from BibleCrawler import BibleCrawler
from PresentationBuilder import PresentationBuilder
from pptx import Presentation


def main():
    # python 3에서는 print() 으로 사용합니다.
    print("Start Program")
    pptx_fpath = './220605_청년부예배.pptx'
    presentation = Presentation(pptx_fpath)
    pptBuilder = PresentationBuilder(presentation)

    # 본문
    book = "사사기"  # 입력은 한글로 받는다
    # 장
    chaper = 1
    # 몇절부터
    start_verse = 1
    # 몇절까지
    end_verse = 10
    # 피피티에서 본문 말씀이 시작되는 위치
    body_start_page = 50

    # 크롤러객체
    bibleCrawler = BibleCrawler()
    bible_list = bibleCrawler.get_bible(book, chaper, start_verse, end_verse)

    # 피피티 파일 전체 슬라이드 장수
    before_ppt_slide_size = len(presentation.slides)

    # 본문 말씀 총 절수
    body_page = end_verse - start_verse

    # 본문 말씀 절수만큼 새 슬라이드를 맨뒤에 추가하기
    for i in range(0, body_page):
        pptBuilder.add_slide(base_slide_page=body_start_page)

    # 추가한 후 전체 슬라이드 장수
    after_ppt_slide_size = len(presentation.slides)

    # 추가된 슬라이드들 중 첫번째 슬라이드의 인덱스
    index_of_first_copied_page = after_ppt_slide_size - body_page

    # 추가된 슬라이드들 중 마지막 놈의 인덱스
    index_of_last_copied_page = after_ppt_slide_size - 1

    # 추가된 슬라이드들의 인덱스를 가리키는 포인터
    index_of_current_copied_page = body_start_page

    # 추가된 슬라이드들을 포문으로 돌면서 본문 시작위치 뒤로 이동시키기
    for index in range(index_of_first_copied_page, index_of_last_copied_page + 1):
        # print("인덱스 " + str(index) + "번째 슬라이드가 " + str(index_of_current_copied_page) + "번째 인덱스로 이동")
        pptBuilder.move_slide(index, index_of_current_copied_page)
        index_of_current_copied_page += 1

    # 중간 저장
    pptBuilder.save_ppt()

    # 이동된 ppt 슬라이드에 크롤링한 성경구절을 넣는다
    pptBuilder.insert_text(start_slide_num=body_start_page,
                           count=end_verse - start_verse,
                           bible_list=bible_list,
                           book=book,
                           chapter=chaper,
                           verse=start_verse)

    # 저장
    pptBuilder.save_ppt()
    print("생성완료")

if __name__ == "__main__":
    main()
