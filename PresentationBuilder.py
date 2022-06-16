import copy
from pptx.enum.text import PP_ALIGN  # 정렬 설정하기
from pptx.util import Pt  # Pt 폰트사이즈
from pptx.dml.color import RGBColor


class PresentationBuilder:
    def __init__(self, resourcePpt):
        self.presentation = resourcePpt

    @property
    def xml_slides(self):
        return self.presentation.slides._sldIdLst

    def add_slide(self, base_slide_page):
        slide_layout = self.presentation.slide_layouts[6]  # 빈 슬라이드 선언
        copy_slide = self.presentation.slides.add_slide(slide_layout)
        base_slide = self.presentation.slides[base_slide_page]
        # 기존 슬라이드에 있던 shape들 복사
        for shape in base_slide.shapes:
            el = shape.element
            newel = copy.deepcopy(el)
            copy_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    def move_slide(self, target, index):
        slides = list(self.xml_slides)
        self.xml_slides.insert(index, slides[target])

    # also works for deleting slides
    def delete_slide(self, index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[index])

    # also works for deleting slides
    def save_ppt(self):
        self.presentation.save('./220605_청년부예배.pptx')

    def insert_text(self, start_slide_num, count, bible_list, book, chapter, verse):
        index = 0
        for i in range(start_slide_num, start_slide_num + count + 1):
            slide_num = start_slide_num + index
            # print("{}번째 반복이다. {}인덱스 슬라이드에 작업중이다".format(index + 1, slide_num))
            slide = self.presentation.slides[slide_num]
            shapes_list = slide.shapes
            shape_index = {}

            # 페이지에 존재하는 shape들의 list
            for i, shape in enumerate(shapes_list):
                shape_index[shape.name] = i

            # shape 이름
            body_contents = "본문내용"
            book_name = "책이름"
            chapter_verse = "몇장몇절"
            rectangle = "직사각형"

            # shape 객체
            body_contents_shape_select = shapes_list[shape_index[body_contents]]
            book_name_shape_select = shapes_list[shape_index[book_name]]
            chapter_verse_shape_select = shapes_list[shape_index[chapter_verse]]
            rectangle_shape_select = shapes_list[shape_index[rectangle]]

            # shape 내 텍스트 프레임 선택하기 & 기존 값 삭제하기
            body_contents_text_frame = body_contents_shape_select.text_frame
            body_contents_text_frame.clear()

            book_name_text_frame = book_name_shape_select.text_frame
            book_name_text_frame.clear()

            chapter_verse_text_frame = chapter_verse_shape_select.text_frame
            chapter_verse_text_frame.clear()

            # 문단 선택하기
            body_contents_p = body_contents_text_frame.paragraphs[0]
            book_name_p = book_name_text_frame.paragraphs[0]
            chapter_verse_p = chapter_verse_text_frame.paragraphs[0]

            # 정렬 설정 -> 중간정렬
            body_contents_p.alighnment = PP_ALIGN.CENTER
            book_name_p.alighnment = PP_ALIGN.CENTER
            chapter_verse_p.alighnment = PP_ALIGN.CENTER

            # 텍스트 입력 / 폰트 지정
            run1 = body_contents_p.add_run()
            run1.text = bible_list[index]
            font1 = run1.font
            font1.name = "맑은 고딕"
            font1.size = Pt(28)
            font1.color.rgb = RGBColor(112, 148, 138)
            font1.bold = True
            # font.name = None  # 지정하지 않으면 기본 글씨체로  #  'Arial'

            run2 = book_name_p.add_run()
            run2.text = book
            font2 = run2.font
            font2.name = "맑은 고딕"
            font2.size = Pt(20)
            font2.color.rgb = RGBColor(242, 119, 113)
            font2.bold = True

            run3 = chapter_verse_p.add_run()
            input_chapter = str(chapter)
            input_verse = str(verse + index)
            run3.text = input_chapter + "장 " + input_verse + "절"
            font3 = run3.font
            font3.name = "맑은 고딕"
            font3.size = Pt(28)
            font3.color.rgb = RGBColor(112, 148, 138)
            font3.bold = True

            index += 1
