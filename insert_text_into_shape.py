from pptx import Presentation  # 라이브러리
from pptx.util import Inches  # 사진, 표등을 그리기 위해
from pptx.enum.text import PP_ALIGN  # 정렬 설정하기
from pptx.util import Pt  # Pt 폰트사이즈

#
# prs = Presentation()  # 파워포인트 객체 선언
#
# for i in range(0, 11):
#     title_slide_layout = prs.slide_layouts[i]  # 슬라이드 종류 선택
#     slide = prs.slides.add_slide(title_slide_layout)  # 슬라이드 추가
#
# prs.save('add all slides.pptx')


# from pptx import Presentation
#
# prs = Presentation()
#
# slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(slide_layout)
#
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "Hello, World!"
# subtitle.text = "python - pptx was here!"
#
#
# prs.save('test.pptx')


# pptx 파일 열기
pptx_fpath = './copy.pptx'
prs = Presentation(pptx_fpath)

# 슬라이드 지정하기
slide_num = 0
slide = prs.slides[slide_num]

# 슬라이드 내 shape 사전 만들기
shapes_list = slide.shapes
shape_index = {}
for i, shape in enumerate(shapes_list):
    shape_index[shape.name] = i
print(shape_index)  # {'Box_down': 0, 'Box_up': 1, 'name2': 2, 'name1': 3}


def text_on_shape(shape, input_text, font_size=95, bold=True):
    # shape 내 텍스트 프레임 선택하기 & 기존 값 삭제하기
    text_frame = shape.text_frame
    text_frame.clear()

    # 문단 선택하기
    p = text_frame.paragraphs[0]

    # 정렬 설정 : 중간정렬
    p.alighnment = PP_ALIGN.CENTER

    # 텍스트 입력 / 폰트 지정
    run = p.add_run()
    run.text = input_text
    font = run.font
    font.name = "궁서체"
    font.size = Pt(font_size)
    font.bold = bold
    #font.name = None  # 지정하지 않으면 기본 글씨체로  #  'Arial'


# 하나의 shape 선택하기
shape_name = '부제목 2'
shape_select = shapes_list[shape_index[shape_name]]

text_on_shape(shape_select, '홍길동')


save_file = './copy.pptx'
prs.save(save_file)