import copy

from pptx import Presentation  # 라이브러리
from pptx.util import Inches  # 사진, 표등을 그리기 위해
from pptx.enum.text import PP_ALIGN  # 정렬 설정하기
from pptx.util import Pt  # Pt 폰트사이즈

pptx_fpath = './220605_청년부예배.pptx'
prs = Presentation(pptx_fpath)

# 슬라이드 지정하기
source_slide = prs.slides[50]

# slide_layouts[6]은 그냥 빈 슬라이드
# [0]은 제목+부제목있는 슬라이드
# 슬라이드 레이아웃 종류는 다음 링크 참고
# https://ggondae.tistory.com/39
slide_layout = prs.slide_layouts[6]  # 빈 슬라이드 선언

copy_slide = prs.slides.add_slide(slide_layout)  # 빈 슬라이드 추가


def xml_slides(self):
    return self.presentation.slides._sldIdLst  # pylint: disable=protected-access


def move_slide(self, old_index, new_index):
    slides = list(self.xml_slides)
    self.xml_slides.remove(slides[old_index])
    self.xml_slides.insert(new_index, slides[old_index])


# 기존 슬라이드에 있던 shape들 복사
for shape in source_slide.shapes:
    el = shape.element
    newel = copy.deepcopy(el)
    copy_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

# move_slide(self=copy_slide, old_index=50, new_index=54)  # Bring last slide to 5th position
# copy_slide.move_to(5)
save_file = './220605_청년부예배.pptx'
prs.save(save_file)

#
# index() 메서드를 사용하여 슬라이드 번호에 액세스할 수 있습니다.
#
# slide_id = prs.slides.index(slide)
# 슬라이드 ID로 슬라이드에 액세스하려면 get() 메서드를 사용할 수 있습니다.
#
# slide = prs.slides.get(slide_id)